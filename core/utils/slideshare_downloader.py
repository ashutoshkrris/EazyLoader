from io import BytesIO
import os
import img2pdf
import re
from os import walk
from os.path import join
import requests
from bs4 import BeautifulSoup
import pptx
import shutil


CURRENT = os.path.dirname(__file__)


class SlideShareDownloader:
    """SlideShare Downloader Class"""

    def __init__(self, slideshare_url=None, download_format='pdf'):
        self.download_format = download_format
        self.slideshare_url = slideshare_url

    def get_slide_info(self):
        html = requests.get(self.slideshare_url).content
        soup = BeautifulSoup(html, 'lxml')
        title = soup.find(class_='j-title-breadcrumb').get_text().strip()
        image_url = soup.find(class_='slide-image')['srcset']
        final_img_url = image_url.split(',')[2].replace(
            ' ', '').replace('1024w', '')
        total_slides = soup.find(id='total-slides').get_text().strip()
        metadata = soup.find_all(class_='metadata-item')
        category = soup.find(class_='slideshow-category').get_text().strip()
        date, views = None, None
        if len(metadata) >= 2:
            date, views = metadata[0].get_text(
            ).strip(), metadata[2].get_text().strip()

        return title, final_img_url, total_slides, category, date, views

    def get_file_name(self):
        # get url basename and replace non-alpha with '_'
        file_name = re.sub('[^0-9a-zA-Z]+', '_',
                           self.slideshare_url.split("/")[-1])
        if file_name.strip() == '':
            print(
                "Something wrong to get filename from URL, fallback to result.pdf or result.pptx")
            file_name = f"result.{self.download_format.lower()}"
        else:
            file_name += f".{self.download_format.lower()}"
        return file_name

    def download_images(self):
        html = requests.get(self.slideshare_url).content
        soup = BeautifulSoup(html, 'lxml')
        # soup.title.string
        title = '/tmp'
        images = soup.findAll('img', {'class': 'slide-image'})
        i = 0
        for image in images:
            image_url = image.get('srcset')
            print(image_url)
            final_img_url = image_url.split(',')[2].replace(
                ' ', '').replace('1024w', '')
            img = requests.get(final_img_url, verify=False)
            if not os.path.exists(title):
                os.makedirs(title)
            with open(f"{title}/{i}", 'wb') as f:
                f.write(img.content)
            i += 1
        print(self.download_format)
        bfr, filename = self.convert(title)
        return bfr, filename

    def convert(self, img_dir_name):
        try:
            imgs = []
            for (dirpath, dirnames, filenames) in walk(join(CURRENT, img_dir_name)):
                imgs.extend(filenames)
                break
            imgs = ["%s/%s" % (img_dir_name, x) for x in imgs]

            def atoi(text):
                return int(text) if text.isdigit() else text

            def natural_keys(text):
                return [atoi(c) for c in re.split('(\d+)', text)]

            imgs.sort(key=natural_keys)

            f_bfr = BytesIO()
            filename = self.get_file_name()
            if self.download_format == 'pdf':
                pdf_bytes = img2pdf.convert(imgs, dpi=300, x=None, y=None)
                with open(filename, "wb") as doc:
                    doc.write(pdf_bytes)

                with open(filename, "rb") as fp:
                    f_bfr.write(fp.read())
                f_bfr.write(pdf_bytes)
            else:
                p = pptx.Presentation()
                blank_slide_layout = p.slide_layouts[6]
                for im in imgs:
                    slide = p.slides.add_slide(blank_slide_layout)
                    slide.shapes.add_picture(
                        im, 0, 0, p.slide_width, p.slide_height)

                p.save(filename)
                with open(filename, "rb") as fp:
                    f_bfr.write(fp.read())

            f_bfr.seek(0)
            os.remove(filename)
            shutil.rmtree(join(CURRENT, img_dir_name))
            return f_bfr, filename

        except Exception as e:
            print(e)
            return None, None
